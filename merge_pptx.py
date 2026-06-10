#!/usr/bin/env python3
"""
merge_pptx.py — Merge TERESA_WBC_v2_27May.pptx with TERESA_comparison.pptx

Logic:
  1. Copy slides 1–41 (0-indexed 0–40) from original → keep
  2. SKIP slide 42 (0-indexed 41, "THANK YOU")
  3. Copy slides 43–44 (0-indexed 42–43) from original
  4. Copy ALL 15 slides from TERESA_comparison.pptx
  5. Append a new "THANK YOU" slide
  6. Output: TERESA_WBC_v3.pptx

Usage:
  python3 merge_pptx.py
"""

from pathlib import Path
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree


# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).resolve().parent
UNIFE_DIR = Path("/Users/andrea/Documents/UNIFE/DOTTORATO/Presentazioni/TERESA")

ORIGINAL = UNIFE_DIR / "TERESA_WBC_v2_27May.pptx"
COMPARISON = UNIFE_DIR / "TERESA_comparison.pptx"
OUTPUT = UNIFE_DIR / "TERESA_WBC_v3.pptx"


def clone_slide(src_slide, dest_prs):
    """Clone a slide from src_slide into dest_prs, preserving all shapes."""
    # Use the same layout as the source slide
    src_layout = src_slide.slide_layout
    # Find matching layout in destination by name
    dest_layout = None
    for layout in dest_prs.slide_layouts:
        if layout.name == src_layout.name:
            dest_layout = layout
            break
    if dest_layout is None:
        # Fallback: use first layout
        dest_layout = dest_prs.slide_layouts[0]

    new_slide = dest_prs.slides.add_slide(dest_layout)

    # Remove placeholder shapes from the new slide (we'll copy actual shapes)
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Copy all shapes from source
    for shape in src_slide.shapes:
        el = shape._element
        new_el = el.__deepcopy__({})
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def add_thank_you_slide(prs):
    """Add a simple centered 'THANK YOU' slide."""
    # Find blank layout by name, fallback to last layout
    slide_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "empty", "vuoto"):
            slide_layout = layout
            break
    if slide_layout is None:
        slide_layout = prs.slide_layouts[-1]  # fallback to last layout
    slide = prs.slides.add_slide(slide_layout)

    # Remove any placeholders
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Add text box centered on the slide
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(7.0)
    height = Inches(2.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.alignment = PP_ALIGN.CENTER

    return slide


def main():
    # ── Validate inputs ────────────────────────────────────────────────────
    if not ORIGINAL.exists():
        print(f"ERROR: Original file not found: {ORIGINAL}", file=sys.stderr)
        sys.exit(1)

    if not COMPARISON.exists():
        print(f"ERROR: Comparison file not found: {COMPARISON}", file=sys.stderr)
        print("Generate TERESA_comparison.pptx first before running this script.", file=sys.stderr)
        sys.exit(1)

    # ── Load presentations ─────────────────────────────────────────────────
    print(f"Loading original: {ORIGINAL.name}")
    orig_prs = Presentation(str(ORIGINAL))
    orig_count = len(orig_prs.slides)
    print(f"  → {orig_count} slides")

    print(f"Loading comparison: {COMPARISON.name}")
    comp_prs = Presentation(str(COMPARISON))
    comp_count = len(comp_prs.slides)
    print(f"  → {comp_count} slides")

    # ── Create output presentation ─────────────────────────────────────────
    # Use the original as template for layouts/themes
    output_prs = Presentation(str(ORIGINAL))

    # Remove all default slides from the new presentation
    ns = "{http://schemas.openxmlformats.org/package/2006/relationships}"
    while len(output_prs.slides) > 0:
        sldId = output_prs.slides._sldIdLst[0]
        rId = sldId.get(ns + "id") or sldId.get("r:id") or sldId.get("id")
        if rId:
            try:
                output_prs.part.drop_rel(rId)
            except KeyError:
                pass  # rel may not exist, that's okay
        output_prs.slides._sldIdLst.remove(sldId)

    # ── Copy slides 1–41 (0-indexed 0–40) ──────────────────────────────────
    print("\nCopying slides 1–41 from original...")
    for idx in range(0, 41):  # 0-indexed: slides 0–40
        clone_slide(orig_prs.slides[idx], output_prs)
    print(f"  → Copied {41} slides")

    # ── SKIP slide 42 (0-indexed 41, THANK YOU) ────────────────────────────
    print("  → Skipped slide 42 (THANK YOU)")

    # ── Copy slides 43–44 (0-indexed 42–43) ────────────────────────────────
    print("Copying slides 43–44 from original...")
    for idx in range(42, 44):  # 0-indexed: slides 42–43
        clone_slide(orig_prs.slides[idx], output_prs)
    print(f"  → Copied {2} slides")

    # ── Copy all 15 slides from comparison ─────────────────────────────────
    print(f"Copying all {comp_count} slides from comparison...")
    for idx in range(comp_count):
        clone_slide(comp_prs.slides[idx], output_prs)
    print(f"  → Copied {comp_count} slides")

    # ── Add new THANK YOU slide ────────────────────────────────────────────
    print("Adding new THANK YOU slide...")
    add_thank_you_slide(output_prs)

    # ── Save ───────────────────────────────────────────────────────────────
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    output_prs.save(str(OUTPUT))

    total = len(output_prs.slides)
    print(f"\n✅ Merged: {orig_count} + {comp_count} = {total} slides → {OUTPUT.name}")
    print(f"   Saved to: {OUTPUT}")


if __name__ == "__main__":
    main()
